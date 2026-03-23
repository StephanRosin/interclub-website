#!/usr/bin/env python3
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


PPT_PATH = Path.home() / "Downloads" / "TCW_Interclub_Projektuebersicht.pptx"
BACKUP_PATH = Path.home() / "Downloads" / "TCW_Interclub_Projektuebersicht.backup-before-datamodel.pptx"

BLUE = RGBColor(23, 43, 123)
GREEN = RGBColor(95, 166, 55)
TEXT = RGBColor(23, 32, 51)
MUTED = RGBColor(87, 98, 123)
LINE = RGBColor(211, 221, 241)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(250, 252, 255)


def add_header(slide, title: str, subtitle: str = ""):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.3), Inches(12.0), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE

    tx = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(10.8), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Manrope"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE
    if subtitle:
        p2 = tx.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Manrope"
        p2.font.size = Pt(10)
        p2.font.color.rgb = MUTED


def add_table_box(slide, left, top, width, title, fields, accent=BLUE):
    body_h = 0.34 * len(fields) + 0.55
    total_h = 0.55 + body_h
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(total_h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = LINE

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.58))
    header.fill.solid()
    header.fill.fore_color.rgb = accent
    header.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.3), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Manrope"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    y = top + 0.72
    for idx, field in enumerate(fields):
        row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + 0.12), Inches(y), Inches(width - 0.24), Inches(0.28))
        row.fill.solid()
        row.fill.fore_color.rgb = LIGHT if idx % 2 == 0 else WHITE
        row.line.color.rgb = LINE
        rtx = slide.shapes.add_textbox(Inches(left + 0.22), Inches(y + 0.03), Inches(width - 0.44), Inches(0.22))
        rp = rtx.text_frame.paragraphs[0]
        rp.text = field
        rp.font.name = "Manrope"
        rp.font.size = Pt(11.5)
        rp.font.color.rgb = TEXT
        y += 0.32

    return outer


def add_note(slide, left, top, width, text):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = LINE
    tx = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Manrope"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT


def add_relation(slide, x1, y1, x2, y2, label):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = GREEN
    line.line.width = Pt(2.2)
    tx = slide.shapes.add_textbox(Inches((x1 + x2) / 2 - 0.5), Inches((y1 + y2) / 2 - 0.2), Inches(1.2), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = label
    p.font.name = "Manrope"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN


def add_datamodel_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Datenmodell: Teamseite", "Visuelle Darstellung der Tabellen, Spalten und Relationen")

    add_table_box(
        slide, 0.7, 1.55, 3.6, "teams",
        [
            "id (PK)",
            "gender",
            "category",
            "liga",
            "teamziel",
            "trainingstag",
            "created",
        ],
        accent=BLUE,
    )
    add_table_box(
        slide, 4.9, 1.55, 3.8, "players",
        [
            "id (PK)",
            "team_id (FK -> teams.id)",
            "name",
            "captain_status",
            "klassierung",
            "myTennisID",
            "created",
        ],
        accent=GREEN,
    )
    add_table_box(
        slide, 9.3, 1.95, 3.1, "/api/teams JSON",
        [
            "title",
            "teamziel",
            "trainingstag",
            "players[]",
            "sortiert für UI",
        ],
        accent=BLUE,
    )

    add_relation(slide, 4.3, 3.1, 4.9, 3.1, "1:n")
    add_relation(slide, 8.7, 3.1, 9.3, 3.1, "API")

    add_note(slide, 0.8, 5.35, 5.3, "Die Teamseite liest nicht statische HTML-Daten, sondern rendert live aus SQLite über /api/teams.")
    add_note(slide, 6.35, 5.35, 5.2, "Players werden für die UI vorsortiert: Capt., Capt. Stv., danach Klassierung und Name.")


def add_matches_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Datenmodell: Spieltermine", "Von der PDF-Quelle bis zur Anzeige im Frontend")

    add_table_box(
        slide, 0.65, 1.65, 2.7, "SwissTennis PDF",
        [
            "Datum",
            "Zeit",
            "Liga",
            "Runde",
            "Heimteam",
            "Gastteam",
        ],
        accent=BLUE,
    )
    add_table_box(
        slide, 3.85, 1.65, 2.95, "import_clubresult.py",
        [
            "liest PDF",
            "parsed Zeit / Liga",
            "holt Stand-Datum",
            "schreibt JSON",
        ],
        accent=GREEN,
    )
    add_table_box(
        slide, 7.25, 1.45, 2.9, "matches.json",
        [
            "source",
            "updated_at",
            "matches[]",
            "date",
            "time",
            "liga",
            "runde",
            "home",
            "away",
        ],
        accent=BLUE,
    )
    add_table_box(
        slide, 10.45, 1.95, 2.3, "Frontend",
        [
            "fetch matches.json",
            "Runden-Filter",
            "Header-Stand",
        ],
        accent=GREEN,
    )

    add_relation(slide, 3.35, 3.05, 3.85, 3.05, "parse")
    add_relation(slide, 6.8, 3.05, 7.25, 3.05, "write")
    add_relation(slide, 10.15, 3.05, 10.45, 3.05, "read")

    add_note(slide, 0.8, 5.45, 5.2, "Automatisierung auf TrueNAS: täglicher Cronjob um 05:00 lädt das PDF direkt von SwissTennis.")
    add_note(slide, 6.25, 5.45, 6.0, "Der Parser wurde auf rechts-nach-links-Struktur umgestellt, damit Zeiten nicht in der Liga landen.")


def main():
    if not PPT_PATH.exists():
        raise FileNotFoundError(f"Präsentation nicht gefunden: {PPT_PATH}")

    shutil.copy2(PPT_PATH, BACKUP_PATH)
    prs = Presentation(str(PPT_PATH))
    add_datamodel_slide(prs)
    add_matches_slide(prs)
    prs.save(str(PPT_PATH))
    print(PPT_PATH)
    print(BACKUP_PATH)


if __name__ == "__main__":
    main()
