#!/usr/bin/env python3
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT = Path.home() / "Downloads" / "TCW_Interclub_Projektuebersicht.pptx"
BG_IMAGE = BASE_DIR / "bg-home.jpg"

BLUE = RGBColor(23, 43, 123)
GREEN = RGBColor(95, 166, 55)
TEXT = RGBColor(23, 32, 51)
MUTED = RGBColor(87, 98, 123)
LINE = RGBColor(211, 221, 241)
WHITE = RGBColor(255, 255, 255)


def add_full_bg(slide, image_path: Path, prs: Presentation):
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 18
    overlay.line.fill.background()


def add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG_IMAGE, prs)

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.0), Inches(8.6), Inches(3.0))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WHITE
    banner.line.color.rgb = LINE

    tx = slide.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(7.8), Inches(1.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Manrope"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE

    p2 = tx.text_frame.add_paragraph()
    p2.text = subtitle
    p2.font.name = "Manrope"
    p2.font.size = Pt(15)
    p2.font.color.rgb = MUTED


def add_header(slide, title: str, subtitle: str = ""):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.3), Inches(12.0), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE

    tx = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(10.5), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Manrope"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE
    if subtitle:
        p2 = tx.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Manrope"
        p2.font.size = Pt(10)
        p2.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.8, top=1.55, width=11.2, height=5.4, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Manrope"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_two_col_bullets(slide, left_title, left_items, right_title, right_items):
    left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(5.7), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = WHITE
    left_box.line.color.rgb = LINE
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.45), Inches(5.7), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = WHITE
    right_box.line.color.rgb = LINE

    for title, items, left in [(left_title, left_items, 0.85), (right_title, right_items, 6.75)]:
        tx = slide.shapes.add_textbox(Inches(left), Inches(1.7), Inches(4.8), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Manrope"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE
        add_bullets(slide, items, left=left, top=2.1, width=4.9, height=4.5, font_size=15)


def add_timeline(slide, steps):
    x = 0.8
    for idx, (title, body) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.25), Inches(3.8))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE

        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.15), Inches(1.95), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = str(idx + 1)
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.paragraphs[0].font.name = "Manrope"
        badge_tf.paragraphs[0].font.size = Pt(12)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = WHITE

        tx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.45), Inches(1.85), Inches(2.9))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Manrope"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_after = Pt(8)
        p2 = tx.text_frame.add_paragraph()
        p2.text = body
        p2.font.name = "Manrope"
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT

        x += 2.9


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "TCW Interclub 2026",
        "Projektübersicht Website, Datenquellen, technische Umsetzung und Automatisierung auf TrueNAS",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Ausgangslage", "Was zu Beginn vorhanden war")
    add_bullets(slide, [
        "Bereits vorhandene Team-HTML im Downloads-Ordner als Design- und Inhaltsbasis.",
        "Trainingsplan nur als Excel-Datei vorhanden.",
        "Spieltermine nur als PDF ClubResult.pdf vorhanden.",
        "Ziel: eine einheitliche Website im TC-Waidberg-Look, gehostet auf TrueNAS.",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Datenquellen", "Welche Informationen woher kommen")
    add_two_col_bullets(
        slide,
        "Statische / importierte Quellen",
        [
            "Trainingsplan aus IC-Trainingsplan.xlsx.",
            "Spieltermine aus ClubResult.pdf bzw. später direkt von SwissTennis.",
            "Hintergrundbild aus Home_Drohne_TCW_DJI_0039.jpg.",
        ],
        "Dynamische Quellen",
        [
            "Teams und Spieler in SQLite auf TrueNAS.",
            "Spieler-Metadaten aus mytennis.ch (URL + Klassierung).",
            "Teamseite liest live aus /api/teams, Spieltermine aus matches.json.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Gebautes Ergebnis", "Was die Lösung heute umfasst")
    add_bullets(slide, [
        "Eine Ein-Seiten-Website mit drei Bereichen: Teams, Trainingsplan, Spieltermine.",
        "Einheitliches Design mit TCW-Farben, Logo, Banner und Hintergrundbild.",
        "Admin-Oberfläche für Pflege von Teams und Spielern auf separatem Port.",
        "Serverbetrieb auf TrueNAS mit Docker-Containern für Site und Admin.",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Technische Architektur", "Frontend, Backend und Betrieb")
    add_two_col_bullets(
        slide,
        "Frontend",
        [
            "index.html als Single-Page-Oberfläche mit JS-basiertem Umschalten.",
            "Teams via Fetch aus /api/teams.",
            "Spieltermine via Fetch aus matches.json.",
        ],
        "Backend / Hosting",
        [
            "SQLite-Datenbank ic_teams.sqlite auf TrueNAS.",
            "Python HTTP-Server für Website und Python Admin-App für CRUD.",
            "Deployment in Docker; Routing über Nginx Proxy Manager.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Datenmodell und Automatisierung", "Wie Inhalte gepflegt und angereichert werden")
    add_bullets(slide, [
        "SQLite mit Tabellen für Teams und Players; Spieler sind mit Teams verknüpft.",
        "Players enthalten Name, Captain-Status, Klassierung und myTennisID.",
        "Beim Anlegen oder Umbenennen werden Klassierung und Spieler-URL automatisch ermittelt.",
        "Aktueller Datenstand auf dem Server: 14 Teams, 119 Spieler, 119 myTennis-Links, 119 Klassierungen.",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Challenges", "Wichtige technische Hürden und Lösungen")
    add_bullets(slide, [
        "Namenssuche: Sonderzeichen, Mehrfach-Vornamen und Apostrophe führten zu Fehlmatches.",
        "Lösung: Normalisierung, Fallback-Suchen und spezielle Varianten wie O'Driscoll -> Driscoll.",
        "PDF-Parsing: Zeiten wurden anfangs teils in die Liga gezogen.",
        "Lösung: Parser auf feste Struktur von rechts nach links umgestellt (Runde / Heim / Gast, danach Datum / Zeit / Liga).",
        "TrueNAS-Host ist eingeschränkt: kein apt, kein pip, kein pdftotext auf dem Host.",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Server-Automation", "Wie die Spieltermine jetzt automatisch aktualisiert werden")
    add_timeline(slide, [
        ("SwissTennis PDF", "Direkter Download von der öffentlichen ClubResult.pdf-URL ohne Login."),
        ("Server Script", "update_matches_server.sh lädt das PDF auf TrueNAS in den Website-Ordner."),
        ("Docker One-Shot", "Ein kurzer python:3.12-alpine-Container installiert poppler-utils und erzeugt matches.json."),
        ("Cron 05:00", "Root-Cron auf TrueNAS startet den Prozess täglich um 05:00 automatisch."),
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Fazit", "Nutzen der Lösung")
    add_bullets(slide, [
        "Zentrale, mobil nutzbare Interclub-Website für Teams, Trainingsplan und Spieltermine.",
        "Reduzierter Pflegeaufwand durch Admin-UI, Datenbank und automatische MyTennis-Anreicherung.",
        "Spieltermine laufen jetzt serverseitig automatisiert über SwissTennis + Cronjob.",
        "Die Architektur bleibt leichtgewichtig, transparent und gut wartbar.",
    ])

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
